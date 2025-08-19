import streamlit as st
import tempfile
import fitz  # PyMuPDF
import docx
import os
import json
from pptx import Presentation
from pptx.dml.color import RGBColor
import google.generativeai as genai
import re

# --------------------------
# Configure Gemini
# --------------------------
genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
model = genai.GenerativeModel("gemini-1.5-flash")

# --------------------------
# File extractors
# --------------------------
def extract_text_from_pdf(path):
    text = ""
    with fitz.open(path) as doc:
        for page in doc:
            text += page.get_text()
    return text

def extract_text_from_docx(path):
    doc = docx.Document(path)
    return "\n".join([p.text for p in doc.paragraphs])

def read_file(uploaded_file):
    with tempfile.NamedTemporaryFile(delete=False) as tmp:
        tmp.write(uploaded_file.read())
        tmp_path = tmp.name
    if uploaded_file.name.endswith(".pdf"):
        return extract_text_from_pdf(tmp_path)
    elif uploaded_file.name.endswith(".docx"):
        return extract_text_from_docx(tmp_path)
    else:
        with open(tmp_path, "r", encoding="utf-8") as f:
            return f.read()

# --------------------------
# Gemini Structured Resume
# --------------------------
def generate_structured_resume(resume_text, jd_text):
    prompt = f"""
    You are an AI assistant. Based on the Resume and Job Description,
    return ONLY a valid JSON object with these exact fields, inside triple backticks (```json ... ```):

    {{
        "Candidate Name": "...",
        "Role Name": "...",
        "Professional Summary": "Min 40 words",
        "Education": "Institution/College Name",
        "Certifications": "Comma separated certifications",
        "Skillset": ["Skill1","Skill2","Skill3","Skill4","Skill5"],
        "Specializations": ["Specialization1", "Specialization2"],
        "Experience": ["10-14 bullet points relevant to the JD and remove all client names, min 20, max 25 words"],
        "Subheader1": "Heading text",
        "CVPoints1": ["1 bullet point under subheader1"],
        "Subheader2": "Heading text",
        "CVPoints2": ["1 bullet point under subheader2"],
        "Subheader3": "Heading text",
        "CVPoints3": ["1 bullet point under subheader3"],
        "Subheader4": "Heading text",
        "CVPoints4": ["1 bullet point under subheader4"],
        "Subheader5": "Heading text",
        "CVPoints5": ["1 bullet point under subheader5"],
        "Subheader6": "Heading text",
        "CVPoints6": ["1 bullet point under subheader6"],
        "Subheader7": "Heading text",
        "CVPoints7": ["1 bullet point under subheader7"],
        "Subheader8": "Heading text",
        "CVPoints8": ["1 bullet point under subheader8"]
    }}

    Resume:
    {resume_text}

    Job Description:
    {jd_text}
    """
    response = model.generate_content(prompt)
    match = re.search(r"```json(.*?)```", response.text, re.DOTALL)
    if match:
        return match.group(1).strip()
    return response.text

# --------------------------
# Formatting Helper
# --------------------------
def copy_formatting(source_run, target_run):
    if source_run.font:
        target_run.font.name = source_run.font.name
        target_run.font.size = source_run.font.size
        target_run.font.bold = source_run.font.bold
        target_run.font.italic = source_run.font.italic
        target_run.font.color.rgb = source_run.font.color.rgb or RGBColor(0, 0, 0)

def insert_bullet_point(shape, para, text):
    new_para = shape.text_frame.add_paragraph()
    new_para.level = para.level
    new_para.text = text
    if para.runs:
        source_run = para.runs[0]
        if new_para.runs:
            copy_formatting(source_run, new_para.runs[0])

# --------------------------
# Placeholder Replacer
# --------------------------
def replace_text_in_shape(shape, data):
    if not shape.has_text_frame:
        return

    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            text = run.text

            # Replace inline placeholders
            replacements = {
                "{{Candidate Name}}": data.get("Candidate Name", ""),
                "{{Role Name}}": data.get("Role Name", ""),
                "{{Professional Summary}}": data.get("Professional Summary", ""),
                "{{Education}}": data.get("Education", ""),
                "{{Certifications}}": data.get("Certifications", "")
            }

            for key, val in replacements.items():
                if key in text:
                    run.text = text.replace(key, val)

            # Skillset and Specializations
            for i, skill in enumerate(data.get("Skillset", []), start=1):
                placeholder = f"{{{{Skillset{i}}}}}"
                if placeholder in text:
                    run.text = text.replace(placeholder, skill)

            for i, spec in enumerate(data.get("Specializations", []), start=1):
                placeholder = f"{{{{Specializations{i}}}}}"
                if placeholder in text:
                    run.text = text.replace(placeholder, spec)

            # Subheaders
            for i in range(1, 9):
                subheader_key = f"Subheader{i}"
                placeholder = f"{{{{{subheader_key}}}}}"
                if placeholder in text:
                    run.text = text.replace(placeholder, data.get(subheader_key, ""))

            # CVPointers (plain text, no bullet)
            for i in range(1, 9):
                cv_key = f"CVPoints{i}"
                cv_placeholder = f"{{{{CVPointer{i}}}}}"
                if cv_placeholder in text:
                    points = data.get(cv_key, [])
                    if points:
                        run.text = text.replace(cv_placeholder, points[0])

        # Experience section (plain text)
        if "{{Experience}}" in para.text:
            points = data.get("Experience", [])
            if points:
                para.text = "\n".join(points)  # plain text block


# --------------------------
# PPT Filler
# --------------------------
def fill_ppt(template_file, data, output_path):
    prs = Presentation(template_file)
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_text_in_shape(shape, data)
    prs.save(output_path)

# --------------------------
# Streamlit App with API Key Gate
# --------------------------
st.set_page_config(page_title="Resume Creation Tool", layout="centered")

if "api_key" not in st.session_state:
    st.session_state.api_key = None
if "authenticated" not in st.session_state:
    st.session_state.authenticated = False

# Step 1: API Key Screen
if not st.session_state.authenticated:
    st.title("🔑 Enter Gemini API Key")
    api_key_input = st.text_input("Gemini API Key", type="password")

    if st.button("Validate API Key"):
        try:
            genai.configure(api_key=api_key_input)
            test_model = genai.GenerativeModel("gemini-1.5-flash")
            test_model.generate_content("Hello")  # simple test call
            st.session_state.api_key = api_key_input
            st.session_state.authenticated = True
            st.success("API Key validated successfully! 🚀")
            st.rerun()
        except Exception as e:
            st.error(f"Invalid API Key or API Error: {e}")

# Step 2: Main Resume Tool (only if authenticated)
else:
    st.title("Resume Creation Tool")
    genai.configure(api_key=st.session_state.api_key)
    model = genai.GenerativeModel("gemini-1.5-flash")

    resume_file = st.file_uploader("Upload Resume (PDF/DOCX)", type=["pdf", "docx"])
    jd_text_input = st.text_area("Paste Job Description here", height=200)

    TEMPLATE_PATH = "EY PPT Template.pptx"

    if st.button("Generate Candidate PPT"):
        if resume_file and jd_text_input.strip():
            resume_text = read_file(resume_file)
            jd_text = jd_text_input.strip()

            st.write("Generating structured data...")
            raw = generate_structured_resume(resume_text, jd_text)

            try:
                data = json.loads(raw)
            except Exception:
                st.error("Gemini output not valid JSON.")
                st.text_area("Raw Output", raw)
                data = None

            if data:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_ppt:
                    output_path = tmp_ppt.name
                fill_ppt(TEMPLATE_PATH, data, output_path)

                with open(output_path, "rb") as f:
                    st.download_button("Download Candidate PPT", f, file_name="Candidate_Profile.pptx")
        else:
            st.warning("Please upload Resume and paste JD.")
